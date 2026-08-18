# -*- coding: utf-8 -*-
"""데모용 매뉴얼 PPT 두 개를 만든다 — 가상 시스템, 회사 정보 없음.

    python make_manual_ppt.py

manual-import 앱이 PPT 를 읽는 규칙(app/pptx_parse.py)에 맞춰,
**까다로운 부분이 전부 드러나도록** 짰다.

  · 첫 장 = 표지        그림 없고 제목뿐 → 매뉴얼 이름이 된다
  · 같은 제목 두 장     한 항목으로 묶인다
  · 표                  · 로 펴져서 절차 줄이 된다
  · 「끝」 장           항목이 아니라 문의 안내로 간다
  · 작은 그림           120px 미만은 버려진다 (일부러 하나 넣었다)

그림은 진짜 화면 캡처 대신 **가짜 UI 를 그려서** 넣는다. 캡처를 쓰면
그 안에 시스템 이름·주소·계정이 남는데, 그건 지울 방법이 없다.

python-pptx 로 만든다. 파서가 python-pptx 로 읽기 때문에, 같은 도구로
만들면 「제목 자리(placeholder idx 0)」 같은 것이 어긋날 일이 없다.
"""
from __future__ import annotations

import io
import os

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt

# 한글 글꼴. PIL 의 기본 글꼴에는 한글이 없어서, 지정하지 않으면
# 가짜 화면 안 글자가 전부 네모(▯)로 그려진다.
_FONT_FILES = [
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    "/usr/share/fonts/opentype/noto/NotoSansCJK-DemiLight.ttc",
    "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
    "C:/Windows/Fonts/malgun.ttf",
]


def font(size: int):
    for f in _FONT_FILES:
        if os.path.exists(f):
            try:
                return ImageFont.truetype(f, size)
            except Exception:  # noqa: BLE001
                continue
    return ImageFont.load_default()


F_SM = None   # 아래에서 채운다 (import 시점에 파일을 열지 않으려고)
F_MD = None

OUT = os.path.dirname(os.path.abspath(__file__))

NAVY = (31, 56, 100)
GREY = (110, 118, 132)
LINE = (208, 213, 221)
BG = (247, 248, 250)
WHITE = (255, 255, 255)
BLUE = (37, 99, 235)


# ─────────────────────────────────────────────────────────────
# 가짜 화면 그림
# ─────────────────────────────────────────────────────────────
def ui_window(title: str, rows: list[tuple[str, str]], button: str,
              w: int = 760, h: int = 460, highlight: int = -1) -> bytes:
    """창 하나짜리 가짜 화면. 입력칸 몇 개와 버튼 하나."""
    im = Image.new("RGB", (w, h), BG)
    d = ImageDraw.Draw(im)

    # 창
    x0, y0, x1, y1 = 40, 40, w - 40, h - 40
    d.rectangle([x0, y0, x1, y1], fill=WHITE, outline=LINE, width=2)
    d.rectangle([x0, y0, x1, y0 + 46], fill=NAVY)
    d.text((x0 + 18, y0 + 14), title, fill=WHITE, font=font(17))
    for i, cx in enumerate((x1 - 30, x1 - 56, x1 - 82)):
        d.ellipse([cx - 6, y0 + 17, cx + 6, y0 + 29],
                  fill=(220, 220, 225) if i else (235, 120, 110))

    # 입력칸
    y = y0 + 86
    for i, (label, value) in enumerate(rows):
        d.text((x0 + 34, y + 8), label, fill=GREY, font=font(15))
        bx0, bx1 = x0 + 190, x1 - 40
        on = (i == highlight)
        d.rectangle([bx0, y, bx1, y + 34], fill=WHITE,
                    outline=BLUE if on else LINE, width=2 if on else 1)
        d.text((bx0 + 12, y + 8), value, fill=(30, 30, 35), font=font(15))
        y += 56

    # 버튼
    bw, bh = 130, 40
    bx0, by0 = x1 - 40 - bw, y1 - 30 - bh
    d.rectangle([bx0, by0, bx0 + bw, by0 + bh], fill=NAVY)
    fb = font(16)
    tw = d.textlength(button, font=fb)
    d.text((bx0 + (bw - tw) / 2, by0 + 11), button, fill=WHITE, font=fb)

    buf = io.BytesIO()
    im.save(buf, "PNG")
    return buf.getvalue()


def ui_list(title: str, items: list[str], w: int = 760, h: int = 430) -> bytes:
    """목록이 있는 가짜 화면."""
    im = Image.new("RGB", (w, h), BG)
    d = ImageDraw.Draw(im)
    x0, y0, x1, y1 = 40, 40, w - 40, h - 40
    d.rectangle([x0, y0, x1, y1], fill=WHITE, outline=LINE, width=2)
    d.rectangle([x0, y0, x1, y0 + 46], fill=NAVY)
    d.text((x0 + 18, y0 + 14), title, fill=WHITE, font=font(17))

    # 왼쪽 메뉴
    d.rectangle([x0, y0 + 46, x0 + 170, y1], fill=(245, 246, 249))
    for i, m in enumerate(["홈", "결재함", "게시판", "일정", "설정"]):
        yy = y0 + 70 + i * 36
        if i == 1:
            d.rectangle([x0 + 8, yy - 6, x0 + 162, yy + 22], fill=(226, 234, 250))
        d.text((x0 + 24, yy), m, fill=NAVY if i == 1 else GREY, font=font(15))

    # 목록
    y = y0 + 76
    for it in items:
        d.text((x0 + 200, y), "\u2022", fill=BLUE, font=font(15))
        d.text((x0 + 220, y), it, fill=(30, 30, 35), font=font(15))
        d.line([x0 + 196, y + 26, x1 - 30, y + 26], fill=(238, 240, 244))
        y += 44

    buf = io.BytesIO()
    im.save(buf, "PNG")
    return buf.getvalue()


def tiny_icon() -> bytes:
    """120px 미만 — 파서가 버려야 하는 그림. 일부러 넣는다."""
    im = Image.new("RGB", (64, 64), NAVY)
    d = ImageDraw.Draw(im)
    d.ellipse([16, 16, 48, 48], fill=WHITE)
    buf = io.BytesIO()
    im.save(buf, "PNG")
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────
# 슬라이드 만들기
# ─────────────────────────────────────────────────────────────
def _deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def cover(prs, title: str, sub: str = ""):
    """표지 — 그림 없이 제목만. 파서가 이걸 매뉴얼 이름으로 쓴다."""
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    # 부제는 파서가 「제목과 같은 말뿐인가」를 볼 때 걸리므로 비워 둔다.
    # 뭐라도 적으면 표지가 아니라 본문 장이 되어 버린다.
    if len(s.placeholders) > 1:
        sp = s.placeholders[1]
        sp._element.getparent().remove(sp._element)
    return s


def step_slide(prs, title: str, steps: list[str], image: bytes | None = None,
               extra_small: bool = False):
    """제목 + 절차 + 그림 한 장."""
    s = prs.slides.add_slide(prs.slide_layouts[5])   # 제목만 있는 레이아웃
    s.shapes.title.text = title

    box = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.4), Inches(4.6))
    tf = box.text_frame
    tf.word_wrap = True
    for i, t in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.size = Pt(16)

    if image:
        s.shapes.add_picture(io.BytesIO(image), Inches(6.5), Inches(1.5), width=Inches(6.1))
    if extra_small:
        # 파서가 버려야 하는 작은 그림
        s.shapes.add_picture(io.BytesIO(tiny_icon()), Inches(0.7), Inches(6.3),
                             width=Inches(0.45))
    return s


def table_slide(prs, title: str, header: list[str], rows: list[list[str]]):
    """표가 있는 장 — 파서가 · 로 펴서 절차 줄로 만든다."""
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = title
    shape = s.shapes.add_table(len(rows) + 1, len(header),
                               Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.8))
    tbl = shape.table
    for c, h in enumerate(header):
        tbl.cell(0, c).text = h
    for r, row in enumerate(rows, start=1):
        for c, v in enumerate(row):
            tbl.cell(r, c).text = v
    return s


def closing(prs, note: str):
    """맺음말 — 항목이 아니라 문의 안내로 간다."""
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "끝"
    box = s.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.6))
    p = box.text_frame.paragraphs[0]
    p.text = note
    p.font.size = Pt(18)
    return s


# ─────────────────────────────────────────────────────────────
# 1. 샘플 ERP 접속 매뉴얼
# ─────────────────────────────────────────────────────────────
def make_erp(path):
    prs = _deck()
    cover(prs, "샘플 ERP 접속 매뉴얼")

    # 같은 제목 두 장 → 한 항목으로 묶인다
    step_slide(prs, "1. 프로그램 설치",
               ["사내 자료실에서 설치 파일을 내려받습니다.",
                "내려받은 파일을 오른쪽 버튼으로 눌러 [관리자 권한으로 실행] 을 고릅니다.",
                "설치 경로는 바꾸지 말고 [다음] 을 눌러 진행합니다."],
               ui_window("샘플 ERP 설치", [("설치 경로", "C:\\SampleERP"),
                                           ("바로가기", "바탕화면에 만들기")],
                         "다음", highlight=0),
               extra_small=True)

    step_slide(prs, "1. 프로그램 설치",
               ["설치가 끝나면 바탕화면에 아이콘이 생깁니다.",
                "처음 실행하면 보안 경고가 뜨는데 [허용] 을 누릅니다."],
               ui_window("설치 완료", [("상태", "설치를 마쳤습니다"),
                                       ("판", "3.2.1")], "닫기"))

    step_slide(prs, "2. 접속 정보 입력",
               ["아이디는 사번, 비밀번호는 처음에 받은 임시 비밀번호를 넣습니다.",
                "접속 서버는 목록에서 [운영] 을 고릅니다.",
                "처음 접속하면 비밀번호를 바꾸라는 창이 뜹니다."],
               ui_window("로그인", [("아이디", "20260101"),
                                    ("비밀번호", "••••••••"),
                                    ("접속 서버", "운영")],
                         "접속", highlight=1))

    table_slide(prs, "3. 안 될 때 확인할 것",
                ["증상", "원인", "해결"],
                [["접속 창이 안 뜸", "프로그램이 덜 설치됨", "지우고 다시 설치"],
                 ["비밀번호가 틀리다고 나옴", "임시 비밀번호 기간 지남", "담당자에게 재발급 요청"],
                 ["접속은 되는데 화면이 비어 있음", "권한이 아직 없음", "담당자에게 권한 신청"]])

    closing(prs, "궁금한 점은 IT 담당자에게 문의해 주세요.")
    prs.save(path)
    return path, len(prs.slides.__iter__.__self__._sldIdLst)


# ─────────────────────────────────────────────────────────────
# 2. 샘플 그룹웨어 사용 매뉴얼
# ─────────────────────────────────────────────────────────────
def make_groupware(path):
    prs = _deck()
    cover(prs, "샘플 그룹웨어 사용 매뉴얼")

    step_slide(prs, "1. 로그인",
               ["주소창에 사내 그룹웨어 주소를 넣습니다.",
                "아이디는 사번, 비밀번호는 메일 비밀번호와 같습니다.",
                "[로그인 유지] 는 공용 PC 에서는 켜지 마세요."],
               ui_window("그룹웨어 로그인", [("아이디", "20260101"),
                                             ("비밀번호", "••••••••"),
                                             ("로그인 유지", "끔")],
                         "로그인", highlight=2))

    step_slide(prs, "2. 결재 올리기",
               ["왼쪽 [결재함] 을 누릅니다.",
                "[새 결재] 를 누르고 양식을 고릅니다.",
                "결재선은 팀장 → 본부장 순으로 넣습니다.",
                "첨부는 20MB 까지 됩니다."],
               ui_list("결재함", ["[대기] 8월 비품 구매 요청",
                                  "[진행] 외부 교육 참가 신청",
                                  "[완료] 7월 근태 정정",
                                  "[반려] 출장비 정산"]))

    step_slide(prs, "3. 자주 쓰는 기능",
               ["게시판 글은 [즐겨찾기] 로 모아 둘 수 있습니다.",
                "일정은 팀 달력과 자동으로 맞춰집니다."],
               ui_list("게시판", ["공지 · 사내 시스템 정기 점검 안내",
                                  "공지 · 하계 휴가 신청 방법",
                                  "자유 · 점심 메뉴 추천"]))

    closing(prs, "그룹웨어 관련 문의는 IT 담당자에게 주세요.")
    prs.save(path)
    return path, 0


if __name__ == "__main__":
    a, _ = make_erp(os.path.join(OUT, "1. 샘플 ERP 접속 매뉴얼.pptx"))
    b, _ = make_groupware(os.path.join(OUT, "2. 샘플 그룹웨어 사용 매뉴얼.pptx"))
    for p in (a, b):
        print("만듦:", os.path.basename(p), "%.0fKB" % (os.path.getsize(p) / 1024))
