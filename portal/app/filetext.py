"""파일에서 글자를 뽑아낸다.

### 왜 필요한가

붙인 파일을 처리할 **앱이 없을 때** 쓰는 길이다.
이력서 분석처럼 전용 서비스가 있으면 그쪽이 낫다. 그 앱은 자기 기준으로
점수를 매기고, 결과를 화면에도 남기고, 담당자가 다시 볼 수 있다.

하지만 청구서 한 장을 받아 "이거 요약해줘" 하는 일에까지 앱을 만들 수는 없다.
그럴 때 포털이 글자만 뽑아서 모델에게 넘긴다.

### 뽑기만 한다

여기서는 **해석하지 않는다.** 표를 정리하거나 항목을 골라내는 일은
모델이 한다. 이 파일은 "PDF 안의 글자"를 문자열로 만들어 줄 뿐이다.

### 못 뽑는 것

  - 한글(.hwp) 옛 형식. 뽑아 주는 라이브러리가 마땅치 않다.
  - 옛 오피스 이진 형식(.xls/.doc/.ppt). 새 형식으로 저장하면 된다.
  - 오디오·동영상.

못 뽑으면 **못 뽑는다고 분명히 말한다.** 빈 문자열을 돌려주면
모델이 "내용이 없는 문서"로 착각하고 지어낸다.

### 그림은 여기서 끝나지 않는다

이미지 파일과 스캔한 PDF 는 글자가 아니라 사진이다. 예전에는 여기서
"못 읽습니다" 하고 끝냈지만, 지금은 `NeedsVision` 을 올려서 부르는 쪽이
`vision.py`(그림 읽기)로 넘길 수 있게 한다. 그림 읽기가 꺼져 있으면
그 메시지가 그대로 직원에게 보이므로, 예전 동작과 같아진다.
"""

import csv
import io
import os
import re

# 모델에게 넘길 최대 글자 수. 넘으면 앞부분만 주고 잘렸다고 알린다.
# 너무 크게 잡으면 토큰 비용이 그대로 늘고, 답변 품질도 오히려 떨어진다.
MAX_CHARS = int(os.environ.get("FILE_TEXT_MAX_CHARS", "20000"))

# 글자를 뽑을 수 있는 형식
SUPPORTED = {"pdf", "docx", "xlsx", "pptx", "csv", "txt", "html", "htm", "md", "json"}

# 글자가 아니라 **그림으로 봐야** 하는 형식 (vision.py 가 맡는다)
IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "webp"}


class ExtractError(Exception):
    """사용자에게 그대로 보여줄 수 있는 실패 사유."""


class NeedsVision(ExtractError):
    """글자로는 뽑을 수 없고 **그림으로 봐야** 하는 경우.

    스캔한 PDF 와 이미지 파일이 여기에 해당한다.
    부르는 쪽(tools/file_read.py)이 이것을 받으면 vision.py 로 넘긴다.
    그림 읽기가 꺼져 있으면 이 메시지가 그대로 직원에게 보인다.
    """


def _pdf(blob: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        raise ExtractError("PDF 를 읽는 기능이 설치되어 있지 않습니다. (pypdf)")
    try:
        reader = PdfReader(io.BytesIO(blob))
    except Exception as e:  # noqa: BLE001
        raise ExtractError(f"PDF 를 열지 못했습니다. 손상된 파일일 수 있습니다. ({e})")

    if getattr(reader, "is_encrypted", False):
        # 빈 비밀번호로 열리는 경우가 흔하다. 그것마저 안 되면 포기한다.
        try:
            reader.decrypt("")
        except Exception:  # noqa: BLE001
            raise ExtractError("비밀번호가 걸린 PDF 입니다. 잠금을 풀고 다시 올려주세요.")

    parts = []
    for i, page in enumerate(reader.pages, 1):
        try:
            t = page.extract_text() or ""
        except Exception:  # noqa: BLE001
            t = ""
        if t.strip():
            parts.append(f"[{i}쪽]\n{t.strip()}")
    text = "\n\n".join(parts)

    if not text.strip():
        # 글자가 하나도 없다 = 쪽 전체가 사진이다.
        # 그림으로 읽을 수 있으니 여기서 끝내지 않고 넘긴다.
        raise NeedsVision(
            "이 PDF 에는 글자가 없고 사진만 들어 있습니다 (스캔한 문서). "
            "글자가 들어 있는 PDF 로 다시 올리거나, 내용을 직접 붙여넣어 주세요."
        )
    return text


# 스캔한 PDF 안에 들어 있는 그림을 꺼낼 때 쓸 수 있는 압축 방식.
# JPEG(DCTDecode) 는 스트림 자체가 그대로 .jpg 파일이라 꺼내기만 하면 된다.
# Flate·CCITT·JBIG2 는 이미지 파일로 되살리려면 Pillow 나 별도 처리가 필요한데,
# 그 무게를 지지 않기로 했다. 그런 PDF 는 이유를 말하고 돌려보낸다.
_PDF_IMAGE_FILTERS = {"/DCTDecode": "jpg"}


def pdf_images(blob: bytes, max_pages: int = 5) -> list[tuple[bytes, str]]:
    """스캔한 PDF 에서 쪽별 그림을 꺼낸다. [(바이트, 확장자), ...]

    글자가 없는 PDF 는 대개 쪽마다 사진 한 장이 통째로 들어 있다.
    그 사진을 꺼내 그림 읽기(vision.py)로 넘기기 위한 것이다.

    새 라이브러리를 쓰지 않는다. 이미 있는 pypdf 로 스트림만 꺼낸다.
    """
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(blob))
    if getattr(reader, "is_encrypted", False):
        try:
            reader.decrypt("")
        except Exception:  # noqa: BLE001
            raise ExtractError("비밀번호가 걸린 PDF 입니다. 잠금을 풀고 다시 올려주세요.")

    out: list[tuple[bytes, str]] = []
    for page in reader.pages[:max_pages]:
        try:
            res = page.get("/Resources")
            xobj = res.get_object().get("/XObject") if res else None
            if not xobj:
                continue
            xobj = xobj.get_object()
        except Exception:  # noqa: BLE001
            continue

        for name in xobj:
            try:
                obj = xobj[name].get_object()
                if obj.get("/Subtype") != "/Image":
                    continue
                filt = obj.get("/Filter")
                # 필터가 여러 개 걸린 경우 마지막 것이 그림 압축이다
                if isinstance(filt, list):
                    filt = filt[-1] if filt else None
                ext = _PDF_IMAGE_FILTERS.get(str(filt))
                if not ext:
                    continue
                data = obj.get_data()
                if data:
                    out.append((data, ext))
                    break          # 한 쪽에 한 장이면 충분하다
            except Exception:  # noqa: BLE001
                continue
    return out


# 오피스 새 형식(.xlsx/.docx/.pptx)은 zip 이고, 붙여 넣은 그림은
# 이 폴더 안에 원본 그대로 들어 있다. 꺼내기만 하면 된다.
_OFFICE_MEDIA = {"xlsx": "xl/media/", "docx": "word/media/", "pptx": "ppt/media/"}

# 로고·아이콘·글머리표까지 읽으면 돈만 나간다. 읽을 만한 그림인지는
# **파일 크기가 아니라 화면상의 크기**로 가린다. 잘 눌린 캡처는 20KB 도
# 안 되는 경우가 있고, 화려한 로고는 100KB 가 넘기도 한다.
MIN_IMAGE_PX = int(os.environ.get("DOC_IMAGE_MIN_PX", "200"))

# 크기를 알아내지 못한 형식에만 쓰는 보조 기준
MIN_IMAGE_BYTES = int(os.environ.get("DOC_IMAGE_MIN_BYTES", "10000"))


def _size_of(data: bytes, kind: str) -> tuple[int, int] | None:
    """그림의 가로·세로를 헤더에서 읽는다. 모르면 None.

    Pillow 를 쓰지 않는다. 크기 몇 바이트 읽자고 이미지 라이브러리를
    통째로 넣을 이유가 없다.
    """
    try:
        if kind == "png" and data[:8] == b"\x89PNG\r\n\x1a\n":
            return (int.from_bytes(data[16:20], "big"),
                    int.from_bytes(data[20:24], "big"))
        if kind == "gif" and data[:3] == b"GIF":
            return (int.from_bytes(data[6:8], "little"),
                    int.from_bytes(data[8:10], "little"))
        if kind in ("jpg", "jpeg") and data[:2] == b"\xff\xd8":
            # SOF0~SOF15 표시자 앞에 세로·가로가 들어 있다 (SOF4/8/12 제외)
            i = 2
            while i + 9 < len(data):
                if data[i] != 0xFF:
                    i += 1
                    continue
                marker = data[i + 1]
                if 0xC0 <= marker <= 0xCF and marker not in (0xC4, 0xC8, 0xCC):
                    return (int.from_bytes(data[i + 7:i + 9], "big"),
                            int.from_bytes(data[i + 5:i + 7], "big"))
                i += 2 + int.from_bytes(data[i + 2:i + 4], "big")
    except Exception:  # noqa: BLE001
        return None
    return None


def _worth_reading(data: bytes, kind: str) -> bool:
    wh = _size_of(data, kind)
    if wh:
        return min(wh) >= MIN_IMAGE_PX
    return len(data) >= MIN_IMAGE_BYTES     # 크기를 모르면 용량으로 가린다


def office_images(ext: str, blob: bytes, limit: int = 5) -> tuple[list[tuple[bytes, str]], int]:
    """문서 안에 **붙여 넣은 그림**을 꺼낸다. (그림 목록, 전체 개수)

    엑셀에 성능 측정 화면을 캡처해 붙여 놓는 식의 문서가 사내에 흔하다.
    셀에는 제목뿐이고 숫자는 전부 그림 안에 있다. 글자만 뽑으면
    "헤더만 있고 데이터가 없다"는 엉뚱한 답이 나온다.

    작은 그림(로고·아이콘)은 건너뛴다. 돌려주는 개수는 limit 까지지만,
    전체 개수도 함께 알려준다. 몇 장을 못 봤는지 말해 줘야 하기 때문이다.
    """
    import zipfile

    prefix = _OFFICE_MEDIA.get((ext or "").lower())
    if not prefix:
        return [], 0

    try:
        z = zipfile.ZipFile(io.BytesIO(blob))
    except Exception:  # noqa: BLE001
        return [], 0

    picks: list[tuple[bytes, str]] = []
    total = 0
    with z:
        names = sorted(n for n in z.namelist() if n.startswith(prefix))
        for n in names:
            kind = n.rsplit(".", 1)[-1].lower() if "." in n else ""
            if kind == "jpe":
                kind = "jpg"
            if kind not in IMAGE_EXTS:
                continue
            try:
                data = z.read(n)
            except Exception:  # noqa: BLE001
                continue
            if not _worth_reading(data, kind):
                continue          # 로고·아이콘·글머리표
            total += 1
            if len(picks) < limit:
                picks.append((data, kind))
    return picks, total


def _docx(blob: bytes) -> str:
    try:
        import docx
    except ImportError:
        raise ExtractError("워드 문서를 읽는 기능이 설치되어 있지 않습니다. (python-docx)")
    try:
        d = docx.Document(io.BytesIO(blob))
    except Exception as e:  # noqa: BLE001
        raise ExtractError(f"워드 문서를 열지 못했습니다. ({e})")

    parts = [p.text for p in d.paragraphs if p.text.strip()]
    # 표 안의 내용도 가져온다. 이력서·견적서는 표에 알맹이가 있는 경우가 많다.
    for t in d.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    text = "\n".join(parts)
    if not text.strip():
        raise ExtractError("이 문서에서 글자를 찾지 못했습니다.")
    return text


def _xlsx(blob: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ExtractError("엑셀을 읽는 기능이 설치되어 있지 않습니다. (openpyxl)")
    try:
        wb = load_workbook(io.BytesIO(blob), read_only=True, data_only=True)
    except Exception as e:  # noqa: BLE001
        raise ExtractError(f"엑셀을 열지 못했습니다. 옛 .xls 형식이면 .xlsx 로 저장해 주세요. ({e})")

    out = []
    for ws in wb.worksheets:
        out.append(f"[시트: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                out.append(" | ".join(cells).rstrip(" |"))
            if len("\n".join(out)) > MAX_CHARS:
                break
    wb.close()
    text = "\n".join(out)
    if not text.strip():
        raise ExtractError("이 엑셀에서 내용을 찾지 못했습니다.")
    return text


def _pptx(blob: bytes) -> str:
    """장표에서 글자를 뽑는다.

    슬라이드는 문단이 아니라 **상자**의 모음이다. 그래서 위에서 아래로,
    왼쪽에서 오른쪽으로 정렬해 준다. 그렇게 하지 않으면 제목이 본문 뒤에
    붙어 나오는 일이 생긴다.

    발표자 노트도 가져온다. 장표 본문은 한 줄인데 알맹이는 노트에 있는
    경우가 사내 자료에는 흔하다.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ExtractError("파워포인트를 읽는 기능이 설치되어 있지 않습니다. (python-pptx)")
    try:
        prs = Presentation(io.BytesIO(blob))
    except Exception as e:  # noqa: BLE001
        raise ExtractError(
            f"파워포인트를 열지 못했습니다. 옛 .ppt 형식이면 .pptx 로 저장해 주세요. ({e})"
        )

    out = []
    for i, slide in enumerate(prs.slides, 1):
        lines = []
        shapes = sorted(
            slide.shapes,
            key=lambda s: ((s.top if s.top is not None else 0),
                           (s.left if s.left is not None else 0)),
        )
        for sh in shapes:
            if sh.has_text_frame:
                t = "\n".join(p.text for p in sh.text_frame.paragraphs if p.text.strip())
                if t.strip():
                    lines.append(t.strip())
            if getattr(sh, "has_table", False):
                for row in sh.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        lines.append(" | ".join(cells))
        # 발표자 노트
        try:
            if slide.has_notes_slide:
                n = (slide.notes_slide.notes_text_frame.text or "").strip()
                if n:
                    lines.append(f"(발표자 노트) {n}")
        except Exception:  # noqa: BLE001
            pass
        if lines:
            out.append(f"[{i}번 장표]\n" + "\n".join(lines))
        if len("\n\n".join(out)) > MAX_CHARS:
            break

    text = "\n\n".join(out)
    if not text.strip():
        raise ExtractError(
            "이 파워포인트에서는 글자를 찾지 못했습니다. "
            "장표가 전부 그림으로 되어 있는 경우입니다."
        )
    return text


# 눈에 보이지 않는 부분. 남겨 두면 스타일시트와 스크립트가 통째로 딸려온다.
_HTML_DROP = re.compile(
    r"<(script|style|noscript|template|svg)\b[^>]*>.*?</\1\s*>",
    re.I | re.S,
)
_HTML_BREAK = re.compile(
    r"</(p|div|li|tr|h[1-6]|section|article)\s*>|<br\s*/?>", re.I
)
# 표는 칸을 줄바꿈으로 끊으면 어느 값이 어느 항목인지 알 수 없게 된다.
# 다른 형식(엑셀·워드 표)과 같이 " | " 로 맞춘다.
_HTML_CELL = re.compile(r"</(td|th)\s*>", re.I)


def _html(blob: bytes) -> str:
    """HTML 에서 보이는 글자만 남긴다.

    사내에서 만든 HTML 장표(reveal.js 등)나 저장한 웹페이지가 대상이다.
    파서를 따로 붙이지 않고 표준 라이브러리만 쓴다. 문서를 완벽히 재현하려는
    것이 아니라 **모델이 읽을 글자**를 만드는 것이 목적이다.
    """
    from html import unescape

    raw = _decode(blob)
    raw = _HTML_DROP.sub(" ", raw)
    raw = re.sub(r"<!--.*?-->", " ", raw, flags=re.S)
    # 문단이 끝나는 자리에 줄바꿈을 넣어 둔다. 안 그러면 전부 한 줄이 된다.
    raw = _HTML_CELL.sub(" | ", raw)
    raw = _HTML_BREAK.sub("\n", raw)
    raw = re.sub(r"<[^>]+>", " ", raw)
    text = unescape(raw)
    # 줄 단위로 공백을 정리한다
    lines = [re.sub(r"[ \t ]+", " ", ln).strip() for ln in text.splitlines()]
    lines = [re.sub(r"\s*\|\s*$", "", ln).strip() for ln in lines]
    text = "\n".join(ln for ln in lines if ln)
    if not text.strip():
        raise ExtractError("이 HTML 에서 글자를 찾지 못했습니다.")
    return text


def _csv(blob: bytes) -> str:
    text = _decode(blob)
    try:
        rows = list(csv.reader(io.StringIO(text)))
    except Exception:  # noqa: BLE001
        return text
    return "\n".join(" | ".join(r) for r in rows if any(x.strip() for x in r))


def _decode(blob: bytes) -> str:
    """한글 파일은 UTF-8 이 아닌 경우가 많다. 흔한 순서대로 시도한다."""
    for enc in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
        try:
            return blob.decode(enc)
        except UnicodeDecodeError:
            continue
    return blob.decode("utf-8", errors="replace")


_READERS = {
    "pdf": _pdf,
    "docx": _docx,
    "xlsx": _xlsx,
    "pptx": _pptx,
    "html": _html,
    "htm": _html,
    "csv": _csv,
    "txt": lambda b: _decode(b),
    "md": lambda b: _decode(b),
    "json": lambda b: _decode(b),
}

# 옛 MS 오피스 이진 형식. 읽어 주는 순수 파이썬 라이브러리가 마땅치 않다.
# "왜 안 되는지"와 "그래서 뭘 하면 되는지"를 같이 말해 준다.
_OLD_OFFICE = {
    "xls":  ("옛 엑셀 형식(.xls)", ".xlsx"),
    "doc":  ("옛 워드 형식(.doc)", ".docx"),
    "ppt":  ("옛 파워포인트 형식(.ppt)", ".pptx"),
}


def extract(ext: str, blob: bytes) -> tuple[str, bool]:
    """(글자, 잘렸는지) 를 돌려준다. 못 뽑으면 ExtractError."""
    ext = (ext or "").lower().lstrip(".")

    if ext in _OLD_OFFICE:
        what, better = _OLD_OFFICE[ext]
        raise ExtractError(
            f"{what}은 읽지 못합니다. {better} 로 저장해서 다시 올려주세요."
        )
    if ext in ("hwp", "hwpx"):
        raise ExtractError(
            "한글 문서(.hwp)는 아직 읽지 못합니다. "
            "PDF 로 저장해서 올리시면 읽을 수 있습니다."
        )
    if ext in IMAGE_EXTS:
        # 그림은 여기서 다루지 않는다. 부르는 쪽이 vision.py 로 넘긴다.
        raise NeedsVision(
            "이미지는 글자를 직접 뽑을 수 없습니다. 그림으로 읽어야 합니다."
        )
    if ext not in _READERS:
        raise ExtractError(f"'{ext}' 형식은 읽지 못합니다. "
                           f"가능한 형식: {', '.join(sorted(SUPPORTED))}")

    text = _READERS[ext](blob)
    # 빈 줄이 수십 개씩 들어가는 문서가 많다. 토큰만 먹으므로 정리한다.
    text = re.sub(r"\n{3,}", "\n\n", text).strip()

    if len(text) > MAX_CHARS:
        return text[:MAX_CHARS], True
    return text, False
