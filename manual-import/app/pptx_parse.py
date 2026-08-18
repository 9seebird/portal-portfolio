"""PPT 를 매뉴얼 형태로 읽는다.

### 어떻게 대응시키나

    PPT 파일 하나        →  매뉴얼 하나
    같은 제목의 슬라이드  →  항목(Menu) 하나
    슬라이드 한 장       →  그 항목 안의 묶음 하나 (block)
    슬라이드의 글        →  절차 (steps)
    슬라이드의 그림      →  화면 캡처 (images)

한 가지 일을 설명하는 데 슬라이드가 서너 장 쓰이는 것은 아주 흔하다.
장마다 항목을 하나씩 만들면 목차에 "Menu 4. SAP 사용자 추가" 가 네 번
찍힌다. 그래서 **제목이 같은 슬라이드가 이어지면 한 항목으로 묶는다.**

슬라이드 제목은 **제목 자리(title placeholder)** 를 먼저 본다. PPT 에서
제목 자리에 쓴 글이 곧 그 장의 이름이기 때문이다. 제목 자리를 안 쓰고
그냥 텍스트 상자에 크게 쓴 경우도 흔해서, 그럴 때는 **맨 위에 있는
텍스트 상자의 첫 줄**을 제목으로 본다.

### 표지와 맺음말은 항목으로 세지 않는다

첫 장이 제목뿐이면 표지로 보고, 그 제목을 매뉴얼 이름으로 제안한다.
마지막 장이 "끝" 같은 맺음말이면 항목에서 빼고 **문의 안내**로 돌린다.
둘 다 항목으로 만들면 매뉴얼마다 빈 항목이 앞뒤로 하나씩 생긴다.
"""

from __future__ import annotations

import io
import re
import unicodedata

from pptx import Presentation
from pptx.util import Emu

# 그림이 너무 작으면 화면 장식(아이콘·로고)이다. 캡처가 아니다.
MIN_IMAGE_PX = 120

# 목록 기호는 떼어 낸다. PPT 에서 자동으로 붙는 것이라 글자에 포함되면
# 매뉴얼 화면에서 "• • 설치" 처럼 두 번 찍힌다.
_BULLET = re.compile(r"^[\s·•◦‣▪●○·\-–—*]+")

# PPT 에서 본문 상자 위에 붙여 두는 이름표. 매뉴얼로 옮기면 장마다
# "설명" 이라는 절차가 하나씩 생겨서 눈에 거슬린다.
_LABELS = {"설명", "내용", "절차", "방법", "순서", "description", "steps"}

# 맺음말 장. 여기 적힌 글은 항목이 아니라 문의 안내다.
_CLOSING = {"끝", "이상", "감사합니다", "감사합니다.", "the end", "end", "q&a", "질의응답"}


def _clean(text: str) -> str:
    t = unicodedata.normalize("NFC", text or "")
    t = t.replace("\v", "\n").replace("\x0b", "\n")
    t = _BULLET.sub("", t.strip())
    return re.sub(r"[ \t]+", " ", t).strip()


def _is_title(shape) -> bool:
    try:
        return shape.is_placeholder and shape.placeholder_format.idx == 0
    except Exception:  # noqa: BLE001
        return False


def _top(shape) -> int:
    try:
        return int(shape.top or 0)
    except Exception:  # noqa: BLE001
        return 0


def _texts(shape) -> list[str]:
    """도형 안의 문단들. 표는 줄 단위로 편다."""
    out: list[str] = []
    if getattr(shape, "has_text_frame", False):
        for p in shape.text_frame.paragraphs:
            line = _clean("".join(r.text for r in p.runs) or p.text)
            if line:
                out.append(line)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [_clean(c.text) for c in row.cells]
            cells = [c for c in cells if c]
            if cells:
                out.append(" · ".join(cells))
    return out


def _pictures(shape) -> list:
    """그림 도형. 그룹 안에 든 것도 꺼낸다."""
    out = []
    if shape.shape_type == 6:                    # GROUP
        for sub in shape.shapes:
            out.extend(_pictures(sub))
        return out
    if getattr(shape, "image", None) is not None:
        out.append(shape)
    return out


def _big_enough(pic) -> bool:
    """화면 위 크기로 판단한다. 원본 픽셀은 크지만 아주 작게 넣은 아이콘도 있다."""
    try:
        w = Emu(pic.width).inches * 96
        h = Emu(pic.height).inches * 96
        return min(w, h) >= MIN_IMAGE_PX
    except Exception:  # noqa: BLE001
        return True


def _key(title: str) -> str:
    """묶을 때 쓰는 열쇠. 띄어쓰기와 마침표 차이로 갈리지 않게 한다."""
    t = unicodedata.normalize("NFC", title or "").lower()
    return re.sub(r"[\s.·:\-–—]+", "", t)


def _read_slides(prs) -> list[dict]:
    """슬라이드를 장 단위로 읽는다. 아직 묶지 않는다."""
    slides = []
    for n, slide in enumerate(prs.slides, start=1):
        title = ""
        body: list[str] = []
        pics = []

        shapes = sorted(slide.shapes, key=_top)
        for sh in shapes:
            pics.extend(p for p in _pictures(sh) if _big_enough(p))

        # ① 제목 자리를 먼저 본다
        for sh in shapes:
            if _is_title(sh):
                lines = _texts(sh)
                if lines:
                    title = lines[0]
                break

        # ② 나머지 글을 모은다 (제목으로 쓴 줄과 이름표는 뺀다)
        for sh in shapes:
            if _is_title(sh):
                continue
            for line in _texts(sh):
                if not line or line == title:
                    continue
                if line.strip().strip(":：").lower() in _LABELS:
                    continue
                body.append(line)

        # ③ 제목 자리를 안 썼으면 맨 위 글의 첫 줄을 제목으로
        if not title and body:
            title = body.pop(0)

        images = []
        for p in pics:
            img = p.image
            images.append({"ext": (img.ext or "png").lower(), "bytes": img.blob})

        slides.append({"no": n, "title": title or f"슬라이드 {n}",
                       "steps": body, "images": images})
    return slides


def parse(blob: bytes) -> dict:
    """PPT 를 읽어 매뉴얼 한 벌로 돌려준다. 파일은 저장하지 않는다.

    돌려주는 모양:
        {"title": "표지 제목", "contactNote": "문의 안내",
         "sections": [{"no": 1, "title": "...", "slides": [2, 3],
                       "blocks": [{"no": 2, "steps": [...],
                                   "images": [{"ext","bytes"}]}]}]}
    """
    slides = _read_slides(Presentation(io.BytesIO(blob)))

    # 표지 — 첫 장에 그림이 없고, 본문이 없거나 제목과 같은 말뿐일 때
    cover = ""
    if slides and not slides[0]["images"]:
        first = slides[0]
        rest = [s for s in first["steps"] if _key(s) != _key(first["title"])]
        if not rest:
            cover = first["title"]
            slides = slides[1:]

    # 맺음말 — 마지막 장에 그림이 없고 제목이 맺음말일 때
    note = ""
    if slides and not slides[-1]["images"] and _key(slides[-1]["title"]) in {
            _key(x) for x in _CLOSING}:
        note = " ".join(slides[-1]["steps"]).strip()
        slides = slides[:-1]

    # 제목이 같은 장들을 한 항목으로 묶는다
    sections: list[dict] = []
    for s in slides:
        block = {"no": s["no"], "steps": s["steps"], "images": s["images"]}
        if sections and _key(sections[-1]["title"]) == _key(s["title"]):
            sections[-1]["blocks"].append(block)
            sections[-1]["slides"].append(s["no"])
            continue
        sections.append({"no": len(sections) + 1, "title": s["title"],
                         "slides": [s["no"]], "blocks": [block]})

    return {"title": cover, "contactNote": note, "sections": sections}


# 파일 이름 앞에 붙은 순번. "1. SAP_GUI_매뉴얼" 을 그대로 쓰면 id 가
# 1-sap-gui 가 되어, 이미 있는 sap-gui 를 바꾸는 대신 하나 더 만든다.
_LEADING_NO = re.compile(r"^\s*\d+\s*[.)\-_]\s*")


def slugify(name: str) -> str:
    """파일 이름에서 매뉴얼 id 를 만든다. 영문 소문자·숫자·하이픈만.

    우리말 이름이면 남는 글자가 없으므로, 그때는 빈 값을 돌려주고
    화면에서 사람이 직접 정하게 한다. 자동으로 지어내면 나중에
    무슨 뜻인지 아무도 모르는 id 가 남는다.
    """
    s = unicodedata.normalize("NFKD", name or "").lower()
    s = re.sub(r"\.pptx?m?$", "", s)
    s = _LEADING_NO.sub("", s)
    s = re.sub(r"[\s_]+", "-", s)
    s = re.sub(r"[^a-z0-9-]", "", s)
    s = re.sub(r"-+", "-", s).strip("-")
    return s
